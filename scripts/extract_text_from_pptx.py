import os
from pptx import Presentation

def extraer_texto_de_pptx(ruta_pptx):
    """Extrae el texto de todas las diapositivas de una presentación PowerPoint."""
    presentacion = Presentation(ruta_pptx)
    textos = []

    for diapositiva in presentacion.slides:
        for forma in diapositiva.shapes:
            if hasattr(forma, "text"):
                textos.append(forma.text)

    return "\n".join(textos)

def guardar_texto_en_archivo(texto, ruta_archivo):
    """Guarda el texto en un archivo de texto."""
    with open(ruta_archivo, 'w', encoding='utf-8') as archivo:
        archivo.write(texto)

def extraer_texto_de_todas_las_pptx_en_directorio(directorio):
    """Extrae el texto de todas las presentaciones PowerPoint en un directorio y las guarda en archivos de texto."""
    for nombre_archivo in os.listdir(directorio):
        if nombre_archivo.endswith(".pptx"):
            ruta_pptx = os.path.join(directorio, nombre_archivo)
            texto = extraer_texto_de_pptx(ruta_pptx)
            nombre_txt = os.path.splitext(nombre_archivo)[0] + ".txt"
            ruta_txt = os.path.join(directorio, nombre_txt)
            guardar_texto_en_archivo(texto, ruta_txt)
            print(f"Texto extraído de {nombre_archivo} y guardado en {nombre_txt}")

if __name__ == "__main__":
    directorio = os.getcwd()  # Obtener el directorio actual
    extraer_texto_de_todas_las_pptx_en_directorio(directorio)

